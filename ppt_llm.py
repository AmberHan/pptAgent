import os
import random
from copy import deepcopy

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

from test_data import ppt_json1


def load_presentation(ppt_path):
    return Presentation(ppt_path)


def is_all_digits(text):
    return text.strip().replace(".", "").isdigit()


def check_text_bold(slide):
    all_bold = True
    for paragraph in slide.text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.font.bold:
                all_bold = False
                break
    return all_bold


def has_text_in_group(group_shape):
    for sub_shape in group_shape.shapes:
        if sub_shape.has_text_frame and sub_shape.text_frame.text.strip():
            return True
        elif sub_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            if has_text_in_group(sub_shape):
                return True
    return False


def get_text_elements(slide):
    elements = []
    top_left_shape = None
    min_distance = float('inf')
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP and has_text_in_group(
                shape) or shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            print("当前页有特殊情况")
            return [top_left_shape]
        if not shape.has_text_frame or (shape.has_text_frame and not shape.text_frame.text.strip()):
            continue
        text_frame = shape.text_frame
        if is_all_digits(text_frame.text) or text_frame.text == '':
            continue
        left = shape.left
        top = shape.top
        distance = (left ** 2 + top ** 2)
        if distance < min_distance:
            min_distance = distance
            top_left_shape = shape
        elements.append(shape)
    if top_left_shape is not None:
        elements.remove(top_left_shape)
        elements.insert(0, top_left_shape)
    return elements


def get_digital_elements(slide):
    elements = []
    for shape in slide.shapes:
        if not shape.has_text_frame or (shape.has_text_frame and not shape.text_frame.text.strip()):
            continue
        text_frame = shape.text_frame
        if is_all_digits(text_frame.text):
            elements.append(shape)
    return elements


def get_contents_elements(slide):
    elements = []
    contents = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP and has_text_in_group(
                shape) or shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            print("当前页有特殊情况")
            # return [contents, []]
        if not shape.has_text_frame or (shape.has_text_frame and not shape.text_frame.text.strip()):
            continue
        text_frame = shape.text_frame
        if is_all_digits(text_frame.text) or text_frame.text == '':
            continue
        if judge_contents(text_frame.text):
            contents = shape
        else:
            elements.append(shape)
    return contents, elements


def judge_contents(contents):
    contents = contents.replace("", "")
    if "目录" in contents or "contents" in contents or "CONTENTS" in contents:
        return True
    return False


def group_texts(elements):
    group_a, group_b = [], []
    title = elements[0]
    for shape in elements[1:]:
        text = shape.text_frame.text
        if '页面标题' in text:
            continue
        elif (check_text_bold(shape) or '段落标题' in text) and '段落内容' not in text:
            group_a.append(shape)
        else:
            group_b.append(shape)
    return title, group_a, group_b


def calculate_distance(shape1, shape2):
    center_x1 = shape1.left + shape1.width / 2
    center_y1 = shape1.top + shape1.height / 2
    center_x2 = shape2.left + shape2.width / 2
    center_y2 = shape2.top + shape2.height / 2
    distance = ((center_x2 - center_x1) ** 2 + (center_y2 - center_y1) ** 2)
    return distance


def find_nearest_pairs(group_a, group_b):
    group_c = []
    while group_a and group_b:
        min_diff = float('inf')
        best_pair = (None, None)
        for a in group_a:
            for b in group_b:
                diff = calculate_distance(a, b)
                if diff < min_diff:
                    min_diff = diff
                    best_pair = (a, b)
        if best_pair[0] is not None and best_pair[1] is not None:
            group_c.append(best_pair)
            group_a.remove(best_pair[0])
            group_b.remove(best_pair[1])
        else:
            break
    return group_c


def update_presentation(group_c, update_text):
    if len(update_text) < len(group_c):
        raise ValueError("update_text列表中的元素数量少于配对数量")

    for i, pair in enumerate(group_c):
        shape_a, shape_b = pair[0], pair[1]
        # 获取新的文本内容
        try:
            new_text_a, new_text_b = update_text[i]
            replace_text_with_style(shape_a, new_text_a)
            replace_text_with_style(shape_b, new_text_b)
        except IndexError:
            print(f"警告：update_text列表中缺少第{i}个配对的新文本")
            continue


def copy_paragraph(text_frame, first_paragraph, text):
    new_paragraph = text_frame.add_paragraph()
    # 复制段落对齐方式
    new_paragraph.alignment = first_paragraph.alignment
    for i, run in enumerate(first_paragraph.runs):
        new_run = new_paragraph.add_run()
        new_run.text = text
        new_run.font.bold = run.font.bold
        new_run.font.italic = run.font.italic
        new_run.font.underline = run.font.underline
        new_run.font.size = run.font.size
        break
    return new_paragraph


def replace_text_with_style(shape, new_text):
    text_frame = shape.text_frame
    for i, original_paragraph in enumerate(text_frame.paragraphs):
        original_paragraph.alignment = PP_ALIGN.LEFT
        for j, run in enumerate(original_paragraph.runs):
            if j == 0:
                leading_spaces = len(run.text) - len(run.text.lstrip(' '))
                run.text = ' ' * leading_spaces + new_text
            else:
                run.text = ''
        break


def replace_text_list_with_style(shape, new_text_list):
    text_frame = shape.text_frame
    add_p = len(new_text_list) - len(text_frame.paragraphs)
    for i, original_paragraph in enumerate(text_frame.paragraphs):
        for j, run in enumerate(original_paragraph.runs):
            if j == 0:
                run.text = new_text_list[i]
            else:
                run.text = ''
        original_paragraph.level = 0

    if add_p > 0:
        paragraph = text_frame.paragraphs[0]
        for item in new_text_list[add_p:]:
            p = copy_paragraph(text_frame, paragraph, item)
            p.level = 0


def copy_slide(prs, slide):
    new_slide = prs.slides.add_slide(prs.slide_layouts[slide.slide_layout.slide_master_id - 1])
    for shape in slide.shapes:
        new_shape = deepcopy(shape)
        new_slide.shapes._spTree.insert_element_before(new_shape.element, 'p:extLst')
    return new_slide


def clear_shape(group):
    for shape in group:
        text_frame = shape.text_frame
        text_frame.clear()


def pair_slide(slide, page):
    if page == 1:
        return get_contents_elements(slide)
    elements = get_text_elements(slide)
    title, group_a, group_b = group_texts(elements)
    if len(group_a) != len(group_b):
        print(f"第{page + 1}页不匹配")
    group_c = find_nearest_pairs(group_a, group_b)
    clear_shape(group_a)
    clear_shape(group_b)
    # print
    if title:
        print(f"第{page + 1}页标题: {title.text}")
    for g in group_c:
        print(f"第{page + 1}页配对结果: {g[0].text}:{g[1].text}")
    print("\n")
    return title, group_c


def get_slides_dic(ppt):
    slides = ppt.slides
    slides_dic = {}
    for page, slide in enumerate(slides):
        title, group_c = pair_slide(slide, page)
        # 首页
        if page == 0:
            slides_dic["topic"] = [(slide, title, group_c)]
        # 目录
        elif page == 1:
            slides_dic["directory"] = [(slide, title, group_c)]
        # 标题提示
        elif page == 2 and len(group_c) == 0:
            slides_dic["directory_1"] = [(slide, title, group_c)]
        # 结尾
        elif page == len(ppt.slides) - 1:
            slides_dic["end"] = [(slide, title, group_c)]
        # 正文
        elif len(group_c) != 0:
            key = str(len(group_c))
            if key not in slides_dic:
                slides_dic[key] = [(slide, title, group_c)]
            slides_dic[key].append((slide, title, group_c))
    return slides_dic


def update_all(ppt, order_list, slides_dic, ppt_json):
    # 更新PPT主题
    _, title, _ = slides_dic["topic"][0]
    replace_text_with_style(title, ppt_json["title"])
    # 更新PPT目录
    _, _, elements = slides_dic["directory"][0]
    sections_json = ppt_json["sections"]
    section_len = len(sections_json)
    for i, element in enumerate(elements):
        if i < section_len:
            replace_text_with_style(element, sections_json[i]["sectionTitle"])
        else:
            element.text_frame.clear()
    # 更新章节
    for section_index, section in enumerate(ppt_json["sections"]):
        update_section(ppt, order_list, slides_dic, section, section_index + 1)


def update_section(ppt, order_list, slides_dic, section_items, section_index):
    # 更新章节页
    slide, title, _ = slides_dic["directory_1"][0]
    replace_text_with_style(title, section_items["sectionTitle"])
    # 更新数字
    number_element = get_digital_elements(slide)
    formatted_number = "{:02d}".format(section_index)
    if len(number_element) == 1:
        replace_text_with_style(number_element[0], formatted_number)
    append_shapes(ppt, slide, order_list)
    # 更新下面章节
    slides_items = section_items["slides"]
    for slide_list in slides_items:
        update_slide(ppt, order_list, slides_dic, slide_list)


def pair_slide_dic(slides_dic, titles):
    if str(titles) not in slides_dic:
        raise Exception("匹配不到合适的PPT")
    # slide, title, group_c = random.choice(slides_dic[titles])
    return random.choice(slides_dic[str(titles)])


def update_slide(ppt, order_list, slides_dic, slides_items):
    slide, title, group_c = pair_slide_dic(slides_dic, len(slides_items))
    replace_text_with_style(title, slides_items["title"])
    contents_items = slides_items["contents"]
    for i, pair in enumerate(group_c):
        shape_title, shape_contents = pair[0], pair[1]
        replace_text_with_style(shape_title, contents_items[i]["title"])
        last_contents = contents_items[i]["contents"]
        # if len(last_contents) == 1:
        #     replace_text_with_style(shape_contents, last_contents[0])
        # else:
        #     replace_text_list_with_style(shape_contents, last_contents)
        replace_text_with_style(shape_contents, "\n".join(last_contents))
    append_shapes(ppt, slide, order_list)


def append_shapes(ppt, source_slide, order_list):
    # 创建一个新幻灯片，使用与原幻灯片相同的布局
    new_slide = ppt.slides.add_slide(source_slide.slide_layout)
    # 清除新幻灯片中的所有默认占位符
    for shape in new_slide.shapes:
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)
    new_slide.background._element = deepcopy(source_slide.background._element)
    for shape in source_slide.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    order_list.append(len(ppt.slides) - 1)


def delete_slide(ppt, slide_index):
    xml_slides = ppt.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[slide_index])


def insert_slide_at_position(ppt, new_slide, position):
    slides = list(ppt.slides._sldIdLst)
    new_slide_id = new_slide.slide_id
    slides.insert(position, new_slide_id)
    ppt.slides._sldIdLst.clear()
    for sldId in slides:
        ppt.slides._sldIdLst.append(sldId)
    return new_slide


def resort_slide(ppt, order_list):
    slides = list(ppt.slides._sldIdLst)
    ppt.slides._sldIdLst.clear()
    for index in order_list:
        ppt.slides._sldIdLst.append(slides[index])


def select_random_ppt(directory):
    # 获取所有PPT文件的列表
    ppt_files = [f for f in os.listdir(directory) if f.endswith(('.ppt', '.pptx'))]

    if not ppt_files:
        print("未找到任何PPT文件")
        return None

    # 随机选择一个PPT文件
    selected_file = random.choice(ppt_files)
    print(f"随机选择的PPT文件是：{selected_file}")
    return os.path.join(directory, selected_file)


def generate(ppt_file, md_json):
    ppt = load_presentation(ppt_file)
    slides_count = len(ppt.slides)
    order_list = [0, 1]

    slides_dic = get_slides_dic(ppt)
    update_all(ppt, order_list, slides_dic, md_json)

    order_list.append(slides_count - 1)
    resort_slide(ppt, order_list)
    ppt.save('./runs/final.pptx')


def test():
    ppt = load_presentation("./ppt_templates/test.pptx")
    page = 11
    title, group_c = pair_slide(ppt.slides[page - 1], page - 1)
    replace_text_with_style(title, f"页面标题{page}")
    update_presentation(group_c, [(f"标题{i}", f"内容{i}" * 10) for i in range(len(group_c))])
    ppt.save('./final.pptx')


if __name__ == '__main__':
    # test()

    ppt_file = select_random_ppt("./ppt_templates")
    # ppt_file = "./ppt_templates/1816068952535851008.pptx"
    generate(ppt_file, ppt_json1)
